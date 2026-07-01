"""
程序化PPT QA — 不依赖视觉，纯代码检查布局和内容
用法: python qa_check.py 答辩PPT.pptx
"""
import sys, re
from pptx import Presentation
from pptx.util import Emu, Pt

pptx_path = sys.argv[1] if len(sys.argv) > 1 else "答辩PPT.pptx"
prs = Presentation(pptx_path)

print(f"🖼 总页数: {len(prs.slides)}")
print(f"📐 尺寸: {prs.slide_width/914400:.2f}\" × {prs.slide_height/914400:.2f}\"")
print()

errors = []
warnings = []

for idx, slide in enumerate(prs.slides, 1):
    texts = []
    has_image = False
    max_box_w = 0
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                t = p.text.strip()
                if t:
                    texts.append(t)
            max_box_w = max(max_box_w, shape.width / 914400)
        if shape.shape_type == 13:  # Picture
            has_image = True

    full = '\n'.join(texts)

    # 双括号残留
    if re.search(r'【【', full) or re.search(r'】】', full):
        errors.append(f"Slide {idx}: 双括号残留 【【/】】")

    # 内容页标题过宽
    if max_box_w > 10.0:
        warnings.append(f"Slide {idx}: 有文本框宽 {max_box_w:.1f}\"，可能 >9.5\" 与校名重叠")

    # 无内容页
    if not texts and not has_image:
        warnings.append(f"Slide {idx}: 完全无内容")

print(f"⚠  警告 ({len(warnings)}):")
for w in warnings:
    print(f"  {w}")

print(f"\n❌ 错误 ({len(errors)}):")
for e in errors:
    print(f"  {e}")

if not errors and not warnings:
    print("  ✅ 无问题")

print(f"\n🔍 内容概览:")
for idx, slide in enumerate(prs.slides, 1):
    first_text = ""
    img_mark = ""
    for shape in slide.shapes:
        if shape.has_text_frame and not first_text:
            p = shape.text_frame.paragraphs[0]
            t = p.text.strip()
            if t:
                first_text = t[:60]
        if shape.shape_type == 13:
            img_mark = " 🖼"
    print(f"  {idx:02d}.{img_mark} {first_text}")
