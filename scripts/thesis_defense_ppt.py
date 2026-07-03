
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os, re
from PIL import Image

# ============================================================
# 颜色系统
# ============================================================
NAVY = RGBColor(0x00, 0x33, 0x99)
RED = RGBColor(0xC0, 0x00, 0x00)
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
LIGHT_BLUE_GRAY = RGBColor(0xE9, 0xED, 0xF4)
SECONDARY_BLUE = RGBColor(0x00, 0x5D, 0xA2)
LIGHT_GRAY = RGBColor(0xDD, 0xDD, 0xDD)
GREEN = RGBColor(0x00, 0x80, 0x00)
ORANGE = RGBColor(0xE6, 0x5C, 0x00)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# ─── 布局位置常量 ────────────────────────────────
TITLE_TOP = 0.38      # 主标题 y，与校徽垂直居中
SUBTITLE_TOP = 1.15   # 副标题栏 y
BODY_TOP = 2.0        # 正文起始 y（有副标题栏时）
BODY_TOP_NO_SUB = 1.5 # 正文起始 y（无副标题栏时）

def remove_textframe(shape):
    txBody = shape._element.find(qn('p:txBody'))
    if txBody is not None:
        shape._element.remove(txBody)

class ThesisDefensePPT:
    def __init__(self, output_path, logo_path=None, name_path=None):
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_WIDTH
        self.prs.slide_height = SLIDE_HEIGHT
        self.output_path = output_path
        self.logo_path = logo_path
        self.name_path = name_path

    def _add_subtitle_bar(self, slide, subtitle, font_size=20, font_color=BLACK, bg_color=LIGHT_BLUE_GRAY):
        """创建自带背景色的副标题栏文本框（自适应高度，支持\\n换行）"""
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(SUBTITLE_TOP), Inches(12.3), Inches(0.8))
        # 设置背景
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = bg_color
        txBox.line.fill.background()
        # 设置文本
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        lines = subtitle.split('\n')
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.alignment = PP_ALIGN.LEFT
            p.line_space = 1.3
            p.space_before = Pt(2)
            p.space_after = Pt(2)
            run = p.runs[0]
            run.font.size = Pt(font_size)
            run.font.name = '微软雅黑'
            run.font.color.rgb = font_color
            run.font.bold = False
        return txBox

    def add_textbox(self, slide, left, top, width, height, text, font_size=18, font_name='微软雅黑', font_color=BLACK, bold=False, align=PP_ALIGN.LEFT, line_space=1.3):
        if not text:
            return slide.shapes.add_textbox(left, top, width, height)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        lines = text.split('\n')
        for i, line in enumerate(lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.alignment = align
            p.line_space = line_space
            run = p.runs[0]
            run.font.size = Pt(font_size)
            run.font.name = font_name
            run.font.color.rgb = font_color
            run.font.bold = bold
        return txBox

    def add_rich_text(self, slide, left, top, width, height, rich_lines, font_size=20, font_name='微软雅黑', line_space=1.3):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        for i, line in enumerate(rich_lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.line_space = line_space
            p.space_after = Pt(8)
            if isinstance(line, str):
                parts = re.split(r'【(.*?)】', line)
                for j, part in enumerate(parts):
                    if j % 2 == 0 and part:
                        run = p.add_run()
                        run.text = part
                        run.font.size = Pt(font_size)
                        run.font.name = font_name
                        run.font.color.rgb = BLACK
                        run.font.bold = False
                    elif j % 2 == 1:
                        run = p.add_run()
                        run.text = part
                        run.font.size = Pt(font_size)
                        run.font.name = font_name
                        run.font.color.rgb = RED
                        run.font.bold = True
            elif isinstance(line, tuple):
                text, color, bold = line
                run = p.add_run()
                run.text = text
                run.font.size = Pt(font_size)
                run.font.name = font_name
                run.font.color.rgb = color
                run.font.bold = bold
        return txBox

    def add_rectangle(self, slide, left, top, width, height, fill_color, no_line=True):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        if no_line:
            shape.line.fill.background()
        return shape

    def add_rounded_rect(self, slide, left, top, width, height, fill_color, radius=0.3):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()
        shape.adjustments[0] = radius
        remove_textframe(shape)
        return shape

    def add_triangle(self, slide, left, top, width, height, fill_color):
        shape = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()
        return shape

    def add_line(self, slide, left, top, width, height, color, line_width=2):
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        line.fill.solid()
        line.fill.fore_color.rgb = color
        line.line.fill.background()
        return line

    def add_picture(self, slide, img_path, left, top, width=None, height=None):
        if not img_path or not os.path.exists(img_path):
            return None
        return slide.shapes.add_picture(img_path, left, top, width=width, height=height)

    def add_picture_constrained(self, slide, img_path, left, top, max_width=None, max_height=None, center_h=False):
        if not img_path or not os.path.exists(img_path):
            return None
        try:
            with Image.open(img_path) as img:
                w, h = img.size
        except:
            return self.add_picture(slide, img_path, left, top, width=max_width, height=max_height)
        aspect = w / h
        if max_width:
            calc_h = max_width / aspect
            calc_w = max_width
        else:
            calc_w = Inches(w / 72)
            calc_h = Inches(h / 72)
        if max_height and calc_h > max_height:
            calc_h = max_height
            calc_w = max_height * aspect
        if center_h:
            left = (SLIDE_WIDTH - calc_w) / 2
        return slide.shapes.add_picture(img_path, left, top, width=calc_w, height=calc_h)

    def add_banner(self, slide, top=Inches(0.15), logo_h=Inches(0.85), name_h=Inches(0.65)):
        if self.logo_path and os.path.exists(self.logo_path):
            self.add_picture(slide, self.logo_path, Inches(0.5), top, height=logo_h)
        if self.name_path and os.path.exists(self.name_path):
            # FIX: 纯浮点运算后转Inches，避免 Length * Length
            name_w = Inches(0.65 * 436 / 165)
            name_x = SLIDE_WIDTH - Inches(0.5) - name_w
            self.add_picture(slide, self.name_path, name_x, top + Inches(0.05), height=name_h)

    # ==================== 页面类型 ====================

    def cover_slide(self, title, eng_title, subtitle, reporter, supervisor, date_str):
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        if self.logo_path and os.path.exists(self.logo_path):
            self.add_picture(slide, self.logo_path, Inches(0.5), Inches(0.15), height=Inches(0.85))
        if self.name_path and os.path.exists(self.name_path):
            name_w = Inches(0.65 * 436 / 165)
            name_x = SLIDE_WIDTH - Inches(0.5) - name_w
            self.add_picture(slide, self.name_path, name_x, Inches(0.2), height=Inches(0.65))
        # 横线在校徽下方
        self.add_line(slide, Inches(0.5), Inches(1.05), Inches(12.3), Inches(0.02), NAVY, 2)
        # 蓝色背景：中文标题
        bg = self.add_rectangle(slide, Inches(0), Inches(2.2), SLIDE_WIDTH, Inches(1.4), NAVY)
        remove_textframe(bg)
        self.add_textbox(slide, Inches(0.5), Inches(2.35), Inches(12.3), Inches(1.0), title, font_size=34, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        # 英文标题（蓝色背景下方，NAVY色）
        if eng_title:
            self.add_textbox(slide, Inches(0.5), Inches(3.85), Inches(12.3), Inches(0.55), eng_title, font_size=22, font_name='Arial', font_color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        # 副标题（居中）
        if subtitle:
            self.add_textbox(slide, Inches(0.5), Inches(4.95), Inches(12.3), Inches(0.5), subtitle, font_size=24, font_color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        # 左下角信息块
        info_x = Inches(0.8)
        info_w = Inches(5.0)
        self.add_textbox(slide, info_x, Inches(5.65), info_w, Inches(0.45), f"汇报人：{reporter}", font_size=22, font_color=NAVY, bold=True, align=PP_ALIGN.LEFT)
        if supervisor:
            self.add_textbox(slide, info_x, Inches(6.1), info_w, Inches(0.45), f"导  师：{supervisor}", font_size=22, font_color=NAVY, bold=True, align=PP_ALIGN.LEFT)
        self.add_textbox(slide, info_x, Inches(6.55), info_w, Inches(0.4), date_str, font_size=18, font_color=NAVY, align=PP_ALIGN.LEFT)
        return slide

    def toc_slide(self, items):
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        self.add_banner(slide)
        # 左侧装饰+标题
        self.add_textbox(slide, Inches(1.0), Inches(1.5), Inches(3.5), Inches(0.7), "目录", font_size=48, font_color=NAVY, bold=True)
        self.add_textbox(slide, Inches(1.0), Inches(2.2), Inches(3.5), Inches(0.5), "CONTENTS", font_size=24, font_color=SECONDARY_BLUE, bold=True)
        self.add_triangle(slide, Inches(0.5), Inches(3.2), Inches(2.8), Inches(2.2), SECONDARY_BLUE)
        self.add_triangle(slide, Inches(1.8), Inches(4.0), Inches(2.6), Inches(1.5), SECONDARY_BLUE)
        self.add_line(slide, Inches(5.0), Inches(1.5), Inches(0.02), Inches(5.0), LIGHT_GRAY, 1)
        # 右侧编号列表（居中偏移，避免太靠左）
        y_start = Inches(1.5)
        list_left = Inches(5.8)
        for i, item in enumerate(items):
            y = y_start + i * Inches(0.85)
            self.add_rounded_rect(slide, list_left, y, Inches(0.55), Inches(0.45), NAVY, 0.3)
            self.add_textbox(slide, list_left, y + Inches(0.02), Inches(0.55), Inches(0.41), f"{i+1}", font_size=18, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
            self.add_textbox(slide, list_left + Inches(0.8), y, Inches(6.0), Inches(0.5), item, font_size=22, font_color=BLACK)
        return slide

    def section_divider_slide(self, chapter_num, chapter_title):
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        self.add_banner(slide)
        self.add_textbox(slide, Inches(1.0), Inches(2.5), Inches(3.0), Inches(2.0), str(chapter_num), font_size=120, font_name='Arial', font_color=NAVY, bold=True)
        self.add_line(slide, Inches(4.0), Inches(2.5), Inches(0.04), Inches(2.5), NAVY, 4)
        self.add_textbox(slide, Inches(4.3), Inches(2.8), Inches(8.0), Inches(1.0), chapter_title, font_size=40, font_color=NAVY, bold=True)
        self.add_textbox(slide, Inches(4.3), Inches(3.8), Inches(8.0), Inches(0.5), f"Chapter {chapter_num}", font_size=20, font_name='Arial', font_color=SECONDARY_BLUE)
        return slide

    def _add_hierarchical_title(self, slide, title):
        """
        渲染层级标题，形如 "1. 选题依据—研究背景—研究意义"
        一级标题 32pt 黑色，次级标题 24pt 红色
        """
        txBox = slide.shapes.add_textbox(Inches(1.5), Inches(TITLE_TOP), Inches(9.5), Inches(0.7))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        # 以—分割（兼容有/无空格）
        parts = re.split(r'\s*—\s*', title, maxsplit=1)
        p = tf.paragraphs[0]
        # 一级标题（黑色）
        run = p.add_run()
        run.text = parts[0]
        run.font.size = Pt(32)
        run.font.name = '微软雅黑'
        run.font.color.rgb = BLACK
        run.font.bold = True
        # 次级标题（红色）
        if len(parts) > 1:
            run = p.add_run()
            run.text = '—' + parts[1]
            run.font.size = Pt(24)
            run.font.name = '微软雅黑'
            run.font.color.rgb = RED
            run.font.bold = True
        return txBox

    def content_slide(self, title, subtitle, rich_lines, body_font_size=20, footnotes=None, body_height=4.8):
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        self.add_banner(slide)
        self._add_hierarchical_title(slide, title)
        body_top = Inches(BODY_TOP)
        if subtitle:
            self._add_subtitle_bar(slide, subtitle)
        else:
            body_top = Inches(BODY_TOP_NO_SUB)
        self.add_rich_text(slide, Inches(0.5), body_top, Inches(12.3), Inches(body_height), rich_lines, font_size=body_font_size)
        if footnotes:
            fn_top = Inches(1.8) + Inches(body_height) + Inches(0.05)
            self.add_textbox(slide, Inches(0.5), fn_top, Inches(12.3), Inches(0.4),
                             footnotes, font_size=14, font_color=BLACK, font_name='微软雅黑')
        return slide

    def image_text_slide(self, title, subtitle, rich_lines, img_path, img_caption=None):
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        self.add_banner(slide)
        self._add_hierarchical_title(slide, title)
        self._add_subtitle_bar(slide, subtitle)
        self.add_rich_text(slide, Inches(0.5), Inches(2.0), Inches(7.8), Inches(4.8), rich_lines, font_size=20)
        if img_path and os.path.exists(img_path):
            pic = self.add_picture_constrained(slide, img_path, Inches(8.5), Inches(2.0), max_width=Inches(4.3), max_height=Inches(4.5))
            if img_caption and pic:
                img_bottom = Inches(2.0) + pic.height + Inches(0.15)
                self.add_textbox(slide, Inches(8.5), img_bottom, Inches(4.3), Inches(0.4), img_caption, font_size=14, font_color=DARK_GRAY, bold=True, align=PP_ALIGN.CENTER)
        return slide

    def two_image_slide(self, title, subtitle, img1_path, img2_path, caption1, caption2):
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        self.add_banner(slide)
        self._add_hierarchical_title(slide, title)
        self._add_subtitle_bar(slide, subtitle)
        img_top = Inches(2.1)
        max_h = Inches(3.8)
        img_w = Inches(5.8)
        pic1 = None
        pic2 = None
        if img1_path and os.path.exists(img1_path):
            pic1 = self.add_picture_constrained(slide, img1_path, Inches(0.5), img_top, max_width=img_w, max_height=max_h)
        if img2_path and os.path.exists(img2_path):
            pic2 = self.add_picture_constrained(slide, img2_path, Inches(7.0), img_top, max_width=img_w, max_height=max_h)
        bottom1 = (img_top + pic1.height) if pic1 else img_top
        bottom2 = (img_top + pic2.height) if pic2 else img_top
        cap_y = max(bottom1, bottom2) + Inches(0.15)
        self.add_textbox(slide, Inches(0.5), cap_y, Inches(5.8), Inches(0.4), caption1, font_size=14, font_color=DARK_GRAY, bold=True, align=PP_ALIGN.CENTER)
        self.add_textbox(slide, Inches(7.0), cap_y, Inches(5.8), Inches(0.4), caption2, font_size=14, font_color=DARK_GRAY, bold=True, align=PP_ALIGN.CENTER)
        return slide

    def bullet_slide(self, title, subtitle, bullets):
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        self.add_banner(slide)
        self._add_hierarchical_title(slide, title)
        self._add_subtitle_bar(slide, subtitle)
        y = Inches(2.0)
        n = len(bullets)
        available = 7.0 - 2.0  # from y start to bottom margin
        spacing = min(Inches(1.0), Inches(available / max(n, 1)))
        box_h = Inches(0.9)
        for i, bullet in enumerate(bullets):
            self.add_rounded_rect(slide, Inches(0.5), y, Inches(0.5), Inches(0.4), NAVY, 0.3)
            self.add_textbox(slide, Inches(0.5), y + Inches(0.02), Inches(0.5), Inches(0.36), str(i+1), font_size=18, font_name='Arial', font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
            self.add_rich_text(slide, Inches(1.2), y - Inches(0.02), Inches(11.5), box_h, [bullet], font_size=20)
            y += spacing
        return slide

    def end_slide(self, main_text="汇报结束", sub_text="谢谢"):
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        self.add_banner(slide, top=Inches(0.15), logo_h=Inches(0.85), name_h=Inches(0.65))
        bg = self.add_rectangle(slide, Inches(0), Inches(2.5), SLIDE_WIDTH, Inches(2.8), NAVY)
        remove_textframe(bg)
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(12.3), Inches(1.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = main_text
        p1.alignment = PP_ALIGN.CENTER
        p1.line_space = 1.2
        run1 = p1.runs[0]
        run1.font.size = Pt(48)
        run1.font.name = '微软雅黑'
        run1.font.color.rgb = WHITE
        run1.font.bold = True
        p2 = tf.add_paragraph()
        p2.text = sub_text
        p2.alignment = PP_ALIGN.CENTER
        p2.line_space = 1.2
        run2 = p2.runs[0]
        run2.font.size = Pt(48)
        run2.font.name = '微软雅黑'
        run2.font.color.rgb = WHITE
        run2.font.bold = True
        return slide

    def full_image_slide(self, title, subtitle, img_path, caption=None):
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        self.add_banner(slide)
        self._add_hierarchical_title(slide, title)
        self._add_subtitle_bar(slide, subtitle)
        if img_path and os.path.exists(img_path):
            pic = self.add_picture_constrained(slide, img_path, Inches(0.5), Inches(2.0), max_width=Inches(12.3), max_height=Inches(4.5), center_h=True)
            if caption and pic:
                img_bottom = Inches(1.8) + pic.height + Inches(0.15)
                self.add_textbox(slide, Inches(0.5), img_bottom, Inches(12.3), Inches(0.4), caption, font_size=14, font_color=DARK_GRAY, bold=True, align=PP_ALIGN.CENTER)
        return slide

    def center_image_slide(self, title, subtitle, img_path, caption=None):
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        self.add_banner(slide)
        self._add_hierarchical_title(slide, title)
        self._add_subtitle_bar(slide, subtitle)
        if img_path and os.path.exists(img_path):
            pic = self.add_picture_constrained(slide, img_path, Inches(1.5), Inches(2.0), max_width=Inches(10.3), max_height=Inches(4.5), center_h=True)
            if caption and pic:
                img_bottom = Inches(1.8) + pic.height + Inches(0.15)
                self.add_textbox(slide, Inches(0.5), img_bottom, Inches(12.3), Inches(0.4), caption, font_size=14, font_color=DARK_GRAY, bold=True, align=PP_ALIGN.CENTER)
        return slide

    def four_image_slide(self, title, subtitle, imgs):
        """
        四图网格页：2×2 排列，每图下方图注。
        imgs = [(path, caption, sub_label), ...]   # sub_label如 "(a)工程样机"
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.add_banner(slide)
        self._add_hierarchical_title(slide, title)
        self._add_subtitle_bar(slide, subtitle)

        cols, rows = 2, 2
        grid_w = Inches(11.5)
        grid_h = Inches(4.0)
        cell_w = grid_w / cols
        cell_h = grid_h / rows
        gap_x = Inches(0.3)
        gap_y = Inches(0.2)
        start_x = Inches(0.9)
        start_y = Inches(2.0)

        for idx, (path, caption, label) in enumerate(imgs):
            r = idx // cols
            c = idx % cols
            cx = start_x + c * (cell_w + gap_x)
            cy = start_y + r * (cell_h + gap_y)
            # 图片约束在单元格内
            max_w = cell_w
            max_h = cell_h - Inches(0.45)
            pic = self.add_picture_constrained(slide, path, cx, cy, max_width=max_w, max_height=max_h)
            if pic:
                # 图注
                cap_y = cy + max_h + Inches(0.05)
                self.add_textbox(slide, cx, cap_y, max_w, Inches(0.35),
                                 label, font_size=12, font_color=DARK_GRAY, bold=True, align=PP_ALIGN.CENTER)
        return slide

    def gantt_slide(self, title, subtitle, phases):
        """
        绘制甘特图时间表。
        phases = [(label, start_month, end_month, color, status_text)]
        month 0 = 2026.01, month 16 = 2027.04
        color: GREEN(已完成), ORANGE(进行中), NAVY(待进行)
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.add_banner(slide)
        self._add_hierarchical_title(slide, title)
        self._add_subtitle_bar(slide, subtitle)

        # 甘特图区域
        chart_left = Inches(3.2)
        chart_top = Inches(2.0)
        chart_w = Inches(9.3)
        chart_h = Inches(4.2)
        month_w = chart_w / 16  # 16个月
        bar_offset = Inches(0.3)  # 进度条下移，避免遮挡日期标签

        # 月份标签
        for i in range(16):
            m = (i % 12) + 1
            y_label = 2026 + (i // 12)
            label = f"{y_label}.{m:02d}"
            x = chart_left + month_w * i
            self.add_textbox(slide, x, chart_top - Inches(0.15), month_w, Inches(0.25),
                             label, font_size=11, font_color=DARK_GRAY, align=PP_ALIGN.CENTER, font_name='Arial')
            # 竖线（从日期下方开始）
            line = self.add_line(slide, x, chart_top + Inches(0.1), Inches(0.005), chart_h, LIGHT_GRAY, 1)

        # 右边界竖线
        self.add_line(slide, chart_left + chart_w, chart_top + Inches(0.1), Inches(0.005), chart_h, LIGHT_GRAY, 1)

        bar_h = Inches(0.55)
        gap = Inches(0.2)
        label_w = Inches(2.8)

        for i, (label, s_mon, e_mon, color, status_text) in enumerate(phases):
            y = chart_top + i * (bar_h + gap) + bar_offset
            # 阶段名称
            self.add_textbox(slide, Inches(0.3), y + Inches(0.05), label_w, bar_h,
                             label, font_size=14, font_color=DARK_GRAY, bold=(i==2))
            # 进度条
            bar_x = chart_left + month_w * s_mon
            bar_w = month_w * (e_mon - s_mon)
            rect = self.add_rectangle(slide, bar_x, y, bar_w, bar_h, fill_color=color)
            # 状态标签
            if status_text:
                self.add_textbox(slide, bar_x + bar_w + Inches(0.1), y + Inches(0.05),
                                 Inches(1.5), bar_h, status_text,
                                 font_size=12, font_color=color, bold=True)
        # 图例
        leg_y = chart_top + len(phases) * (bar_h + gap) + bar_offset + Inches(0.15)
        for j, (c, t) in enumerate([(GREEN, "已完成"), (ORANGE, "进行中"), (NAVY, "待进行")]):
            lx = chart_left + j * Inches(2.5)
            self.add_rectangle(slide, lx, leg_y, Inches(0.25), Inches(0.2), fill_color=c)
            self.add_textbox(slide, lx + Inches(0.35), leg_y - Inches(0.02), Inches(2.0), Inches(0.25),
                             t, font_size=12, font_color=DARK_GRAY)

        # 里程碑标记（在图例下方）
        milestone_y = leg_y + Inches(0.4)
        self.add_textbox(slide, chart_left, milestone_y, Inches(7.0), Inches(0.3),
                         "★ 预计初稿完成：2027年3月    计划定稿：2027年4月",
                         font_size=12, font_color=RED, bold=True, align=PP_ALIGN.LEFT)
        return slide

    def save(self):
        self.prs.save(self.output_path)
        print(f"PPT已保存至: {self.output_path}")
        print(f"总页数: {len(self.prs.slides)}")

# ============================================================
# 主程序
# ============================================================


if __name__ == "__main__":
    # 创建 ThesisDefensePPT 实例并生成幻灯片
    ppt = ThesisDefensePPT("output.pptx")
    ppt.save()
